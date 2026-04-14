#!/usr/bin/env python3
"""
ISP-Agent 项目汇报 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
DARK_BLUE = RGBColor(0x1a, 0x3a, 0x5c)
ACCENT_BLUE = RGBColor(0x2c, 0x5a, 0x8c)
LIGHT_BLUE = RGBColor(0xe8, 0xf4, 0xf8)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x28, 0xa7, 0x45)
ORANGE = RGBColor(0xff, 0x99, 0x00)

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, two_column=False):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部色条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if two_column and isinstance(bullets, dict):
        # 两列布局
        left_bullets = bullets.get('left', [])
        right_bullets = bullets.get('right', [])
        
        # 左列
        txBoxL = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
        tfL = txBoxL.text_frame
        tfL.word_wrap = True
        for i, bullet in enumerate(left_bullets):
            if i == 0:
                p = tfL.paragraphs[0]
            else:
                p = tfL.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.space_after = Pt(12)
        
        # 右列
        txBoxR = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
        tfR = txBoxR.text_frame
        tfR.word_wrap = True
        for i, bullet in enumerate(right_bullets):
            if i == 0:
                p = tfR.paragraphs[0]
            else:
                p = tfR.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
            p.space_after = Pt(12)
    else:
        # 单列布局
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY
            p.space_after = Pt(14)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部色条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 表格
    cols = len(headers)
    table_rows = len(rows) + 1
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * table_rows)
    
    table = slide.shapes.add_table(table_rows, cols, x, y, cx, cy).table
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # 数据行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
    
    return slide

# ========== 创建幻灯片 ==========

# 第1页：封面
add_title_slide(prs, "ISP-Agent 项目汇报", "基于LLM的ISP图像/视频处理助手 | v0.4")

# 第2页：目录
add_content_slide(prs, "目录", [
    "项目背景与痛点",
    "项目目标与价值",
    "核心功能架构",
    "技术实现方案",
    "版本里程碑",
    "后续规划"
])

# 第3页：项目背景
add_content_slide(prs, "项目背景与痛点", {
    'left': [
        "图像质量评估依赖 Imatest 等工具",
        "手动操作繁琐，效率低下",
        "版本对比无法量化差异"
    ],
    'right': [
        "参数调优依赖经验，迭代周期长",
        "经验分散在个人笔记中",
        "测试报告需要手动撰写"
    ]
}, two_column=True)

# 第4页：项目目标
add_content_slide(prs, "项目目标", [
    "自动化图像质量分析（BRISQUE/HDR/色彩/噪声）",
    "ISP版本对比自动化，生成量化对比报告",
    "基于AI分析的智能调参建议",
    "中文自然语言ISP领域问答",
    "测试报告自动生成"
])

# 第5页：核心价值
add_content_slide(prs, "核心价值", [
    "🔄 自动化 - 减少Imatest人工操作，实现自动化批量对比",
    "📊 标准化 - 统一对比流程和评估标准",
    "📈 可追溯 - 记录每次对比结果，形成对比历史",
    "🤖 智能化 - AI分析结果，给出调参建议",
    "📝 自动化 - 自动生成规范测试报告"
])

# 第6页：功能架构
add_content_slide(prs, "核心功能架构", [
    "📊 ISP版本对比框架 - Comp12 RAW自动解析、多维度量化对比",
    "🎯 BRISQUE盲图像质量评分 - 无参考图像质量评估",
    "🔆 HDR动态范围分析 - 动态范围(dB)精确测量",
    "🎨 色彩准确性评估 - ΔE色差计算",
    "📹 视频处理分析 - 抽帧、稳定性检测、内容摘要"
])

# 第7页：技术架构
add_content_slide(prs, "技术实现方案", {
    'left': [
        "Python 3.9+ 本地运行时",
        "OpenCV/Pillow/rawpy 图像处理",
        "FFmpeg 视频编解码",
        "NumPy/SciPy 科学计算"
    ],
    'right': [
        "MiniMax API (默认) / OpenAI / Claude",
        "Docker 容器化部署",
        "REST API 标准化接口",
        "飞书/Telegram 多渠道集成"
    ]
}, two_column=True)

# 第8页：版本里程碑
add_table_slide(prs, "版本里程碑", 
    ["版本", "日期", "主要功能", "状态"],
    [
        ["v0.1", "2026-03-20", "基础框架搭建、ISP Pipeline可视化", "已完成"],
        ["v0.2", "2026-03-25", "RAW处理、HDR增强、车载场景", "已完成"],
        ["v0.3", "2026-03-28", "批量处理、问答引擎、导出管理", "已完成"],
        ["v0.4", "2026-04-08", "BRISQUE、ISO16505合规、调参Advisor", "已完成 ✅"],
        ["v1.0", "规划中", "完整ISP对比框架、RAG知识库", "开发中"]
    ])

# 第9页：核心工具
add_table_slide(prs, "核心工具清单", 
    ["工具名称", "描述", "优先级", "状态"],
    [
        ["brisque.py", "BRISQUE盲图像质量评分", "P0", "✅已完成"],
        ["niqe.py", "NIQE自然图像质量评估", "P0", "✅已完成"],
        ["hdr_processor.py", "HDR动态范围处理", "P0", "✅已完成"],
        ["iso16505.py", "ISO16505标准化评估", "P1", "✅已完成"],
        ["tuning_advisor.py", "AI调参建议引擎", "P1", "✅已完成"],
        ["ai_quality_scorer.py", "综合质量评分融合", "P0", "✅已完成"],
    ])

# 第10页：后续规划
add_content_slide(prs, "后续规划", [
    "v1.0 完善 ISP 对比框架，支持更多RAW格式",
    "构建 ISP 领域 RAG 知识库",
    "支持更多ISP参数自动诊断",
    "完善自动化测试覆盖",
    "扩展视频质量分析能力"
])

# 第11页：结束页
add_title_slide(prs, "谢谢！", "ISP-Agent - 让ISP开发更高效 🚀")

# 保存
output_path = os.path.join(os.path.dirname(__file__), "..", "ISP-Agent项目汇报.pptx")
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
